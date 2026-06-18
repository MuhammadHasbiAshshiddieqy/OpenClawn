"""Tool dokumen: pdf_read (baca) + doc_write (docx/pptx/xlsx/md) + pdf_write (PDF).

Semua library murni-Python tanpa dependency sistem (pypdf, python-docx, python-pptx,
openpyxl, reportlab) — pengecualian dependency yang disetujui owner, lihat CLAUDE.md §7.
"""

from pypdf import PdfReader

from infra.config import CONFIG
from infra.workspace import WorkspaceViolation, resolve_in_workspace
from tools.base import Tool

MAX_PDF_CHARS = CONFIG.tool_max_output

# Format dokumen yang didukung doc_write → ekstensi file.
DOC_FORMATS = {"docx", "pptx", "xlsx", "md"}


class PdfReadTool(Tool):
    name = "pdf_read"
    requires_approval = False

    async def execute(self, input_data: dict, vault, db=None) -> dict:
        path = (input_data.get("path") or "").strip()
        if not path:
            return {"error": "path wajib diisi"}
        try:
            safe = resolve_in_workspace(path, CONFIG.workspace_root)
        except WorkspaceViolation as e:
            return {"error": str(e)}
        if not safe.exists():
            return {"error": f"File tidak ditemukan: {path}"}

        try:
            reader = PdfReader(str(safe))
        except Exception as e:  # pypdf melempar beragam error parsing — tangani semua
            return {"error": f"Gagal membuka PDF: {e}"}

        # Halaman opsional: 1-indexed agar natural bagi user. Default: semua.
        page_arg = input_data.get("page")
        total = len(reader.pages)
        try:
            if page_arg is not None:
                idx = int(page_arg) - 1
                if idx < 0 or idx >= total:
                    return {"error": f"Halaman {page_arg} di luar rentang (1..{total})"}
                pages = [reader.pages[idx]]
            else:
                pages = reader.pages
            text = "\n".join(p.extract_text() or "" for p in pages)
        except Exception as e:
            return {"error": f"Gagal ekstrak teks: {e}"}

        return {
            "pages": total,
            "text": text[:MAX_PDF_CHARS],
            "truncated": len(text) > MAX_PDF_CHARS,
        }

    def schema(self) -> dict:
        return {
            "name": "pdf_read",
            "description": (
                "Ekstrak teks dari file PDF dalam workspace. "
                "Opsional 'page' (1-indexed) untuk satu halaman saja. "
                "Pakai untuk membaca dokumen/laporan/spec berformat PDF."
            ),
            "input_schema": {
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "Path PDF relatif ke workspace."},
                    "page": {
                        "type": "integer",
                        "description": "Nomor halaman (1-indexed). Kosong = semua.",
                    },
                },
                "required": ["path"],
            },
        }


class DocWriteTool(Tool):
    """Tulis dokumen terstruktur (docx/pptx/xlsx/md) ke workspace dari `content`.

    Destruktif (menulis ke filesystem) → requires_approval=True, konsisten file_write.
    Library doc di-import LAZY di dalam execute agar dependency opsional yang hilang
    menghasilkan error anggun, bukan menggagalkan import seluruh modul tools (§1.3).
    """

    name = "doc_write"
    requires_approval = True

    async def execute(self, input_data: dict, vault, db=None) -> dict:
        path = (input_data.get("path") or "").strip()
        fmt = (input_data.get("format") or "").strip().lower()
        content = input_data.get("content")
        if not path:
            return {"error": "path wajib diisi"}
        if fmt not in DOC_FORMATS:
            return {"error": f"format harus salah satu: {', '.join(sorted(DOC_FORMATS))}"}
        try:
            safe = resolve_in_workspace(path, CONFIG.workspace_root)
        except WorkspaceViolation as e:
            return {"error": str(e)}

        try:
            match fmt:
                case "md":
                    self._write_md(safe, content)
                case "docx":
                    self._write_docx(safe, content)
                case "pptx":
                    self._write_pptx(safe, content)
                case "xlsx":
                    self._write_xlsx(safe, content)
        except ImportError as e:
            return {"error": f"Library untuk format {fmt} tidak terpasang: {e}"}
        except (ValueError, TypeError) as e:
            return {"error": f"Struktur content tidak valid untuk {fmt}: {e}"}
        except Exception as e:  # noqa: BLE001 — penulisan dokumen harus gagal anggun
            return {"error": f"Gagal menulis {fmt}: {e}"}

        return {"path": path, "format": fmt, "ok": True}

    # ── Penulis per-format. content fleksibel: str (md) atau struktur (lainnya). ──

    def _write_md(self, safe, content) -> None:
        """Markdown/teks: content boleh string langsung, atau {title, sections:[{heading,body}]}."""
        if isinstance(content, str):
            text = content
        elif isinstance(content, dict):
            parts = []
            if content.get("title"):
                parts.append(f"# {content['title']}\n")
            for sec in content.get("sections", []):
                if sec.get("heading"):
                    parts.append(f"## {sec['heading']}")
                if sec.get("body"):
                    parts.append(str(sec["body"]) + "\n")
            text = "\n".join(parts)
        else:
            raise TypeError("content md harus string atau objek {title, sections}")
        safe.write_text(text, encoding="utf-8")

    def _write_docx(self, safe, content) -> None:
        """Word: content = {title?, sections:[{heading?, body?, bullets?[]}]}."""
        from docx import Document

        if not isinstance(content, dict):
            raise TypeError("content docx harus objek {title, sections}")
        doc = Document()
        if content.get("title"):
            doc.add_heading(str(content["title"]), level=0)
        for sec in content.get("sections", []):
            if sec.get("heading"):
                doc.add_heading(str(sec["heading"]), level=1)
            if sec.get("body"):
                doc.add_paragraph(str(sec["body"]))
            for b in sec.get("bullets", []):
                doc.add_paragraph(str(b), style="List Bullet")
        doc.save(str(safe))

    def _write_pptx(self, safe, content) -> None:
        """PowerPoint: content = {title?, slides:[{title, bullets:[]}]}."""
        from pptx import Presentation

        if not isinstance(content, dict):
            raise TypeError("content pptx harus objek {title, slides}")
        prs = Presentation()
        slides = content.get("slides", [])
        if content.get("title"):  # slide judul
            layout = prs.slide_layouts[0]
            s = prs.slides.add_slide(layout)
            s.shapes.title.text = str(content["title"])
        for sl in slides:
            layout = prs.slide_layouts[1]  # judul + konten
            s = prs.slides.add_slide(layout)
            s.shapes.title.text = str(sl.get("title", ""))
            body = s.placeholders[1].text_frame
            bullets = sl.get("bullets", [])
            for i, b in enumerate(bullets):
                para = body.paragraphs[0] if i == 0 else body.add_paragraph()
                para.text = str(b)
        prs.save(str(safe))

    def _write_xlsx(self, safe, content) -> None:
        """Spreadsheet: content = {sheet?, rows:[[..],[..]]} atau {headers:[], rows:[]}."""
        from openpyxl import Workbook

        if not isinstance(content, dict):
            raise TypeError("content xlsx harus objek {rows} / {headers, rows}")
        wb = Workbook()
        ws = wb.active
        ws.title = str(content.get("sheet", "Sheet1"))[:31]
        if content.get("headers"):
            ws.append([str(h) for h in content["headers"]])
        for row in content.get("rows", []):
            if not isinstance(row, list):
                raise TypeError("setiap rows harus list sel")
            ws.append(row)
        wb.save(str(safe))

    def schema(self) -> dict:
        return {
            "name": "doc_write",
            "description": (
                "Tulis dokumen ke workspace dalam format docx/pptx/xlsx/md. "
                "content menyesuaikan format: md→string atau {title,sections}; "
                "docx→{title,sections:[{heading,body,bullets}]}; "
                "pptx→{title,slides:[{title,bullets}]}; xlsx→{headers,rows:[[..]]}. "
                "Butuh approval (menulis file)."
            ),
            "input_schema": {
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "Path output relatif ke workspace."},
                    "format": {
                        "type": "string",
                        "enum": sorted(DOC_FORMATS),
                        "description": "Format dokumen.",
                    },
                    "content": {
                        "description": "Isi terstruktur sesuai format (lihat description).",
                    },
                },
                "required": ["path", "format", "content"],
            },
        }


class PdfWriteTool(Tool):
    """Tulis PDF terstruktur ke workspace via reportlab (murni-Python).

    Destruktif (menulis file) → requires_approval=True. reportlab di-import lazy
    agar dependency hilang gagal anggun (§1.3). content = {title?, sections:[{heading?,
    body?, bullets?:[]}]} — sama bentuk dengan doc_write docx untuk konsistensi.
    """

    name = "pdf_write"
    requires_approval = True

    async def execute(self, input_data: dict, vault, db=None) -> dict:
        path = (input_data.get("path") or "").strip()
        content = input_data.get("content")
        if not path:
            return {"error": "path wajib diisi"}
        if not isinstance(content, dict):
            return {"error": "content harus objek {title, sections}"}
        try:
            safe = resolve_in_workspace(path, CONFIG.workspace_root)
        except WorkspaceViolation as e:
            return {"error": str(e)}

        try:
            self._render(str(safe), content)
        except ImportError as e:
            return {"error": f"reportlab tidak terpasang: {e}"}
        except Exception as e:  # noqa: BLE001 — penulisan PDF harus gagal anggun
            return {"error": f"Gagal menulis PDF: {e}"}
        return {"path": path, "format": "pdf", "ok": True}

    def _render(self, out_path: str, content: dict) -> None:
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.styles import getSampleStyleSheet
        from reportlab.platypus import ListFlowable, ListItem, Paragraph, SimpleDocTemplate, Spacer

        styles = getSampleStyleSheet()
        flow = []
        if content.get("title"):
            flow.append(Paragraph(str(content["title"]), styles["Title"]))
            flow.append(Spacer(1, 12))
        for sec in content.get("sections", []):
            if sec.get("heading"):
                flow.append(Paragraph(str(sec["heading"]), styles["Heading2"]))
            if sec.get("body"):
                flow.append(Paragraph(str(sec["body"]), styles["BodyText"]))
            bullets = sec.get("bullets", [])
            if bullets:
                items = [ListItem(Paragraph(str(b), styles["BodyText"])) for b in bullets]
                flow.append(ListFlowable(items, bulletType="bullet"))
            flow.append(Spacer(1, 8))
        SimpleDocTemplate(out_path, pagesize=A4).build(flow)

    def schema(self) -> dict:
        return {
            "name": "pdf_write",
            "description": (
                "Tulis dokumen PDF ke workspace dari struktur {title?, sections:[{heading?, "
                "body?, bullets?:[]}]}. Untuk Word pakai doc_write (format docx). "
                "Butuh approval (menulis file)."
            ),
            "input_schema": {
                "type": "object",
                "properties": {
                    "path": {"type": "string", "description": "Path .pdf relatif ke workspace."},
                    "content": {"description": "Struktur {title, sections} (lihat description)."},
                },
                "required": ["path", "content"],
            },
        }
